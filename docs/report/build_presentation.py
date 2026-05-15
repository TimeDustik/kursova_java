"""
Генератор презентації до курсового проекту
«Система управління інвентарем для ІТ-департаменту»
Студент: Скакун М. М., група ТЦР-33
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu
import os

# ─── Кольорова схема ──────────────────────────────────────────────────────────
DARK_BLUE   = RGBColor(0x1A, 0x2A, 0x4A)   # фон заголовка
ACCENT_BLUE = RGBColor(0x0D, 0x6E, 0xFD)   # акцент
LIGHT_BG    = RGBColor(0xF0, 0xF4, 0xFF)   # світлий фон
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT   = RGBColor(0x1A, 0x1A, 0x2E)
GRAY        = RGBColor(0x6C, 0x75, 0x7D)
GREEN       = RGBColor(0x19, 0x87, 0x54)
RED_COLOR   = RGBColor(0xDC, 0x35, 0x45)

OUTPUT = "/Users/mihailskakun/IdeaProjects/kursova_java/docs/report/Презентація_Скакун_ТЦР33.pptx"
UML_CLASS = "/Users/mihailskakun/IdeaProjects/kursova_java/docs/uml-png/IT-Inventory Class Diagram.png"
UML_SEQ   = "/Users/mihailskakun/IdeaProjects/kursova_java/docs/uml-png/Login Sequence.png"

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]   # повністю порожній layout

# ─── Допоміжні функції ────────────────────────────────────────────────────────

def bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def rect(slide, l, t, w, h, color, alpha=None):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def txt(slide, text, l, t, w, h,
        size=20, bold=False, italic=False, color=WHITE,
        align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb

def add_bullet_box(slide, items, l, t, w, h,
                   size=18, color=DARK_TEXT, marker="●", spacing=0.55):
    """Додає блок із маркованим списком."""
    y = t
    for item in items:
        indent = 0
        text = item
        if item.startswith("  "):
            indent = 0.3
            text = item.strip()
            m = "–"
            s = size - 2
        else:
            m = marker
            s = size
        txb = slide.shapes.add_textbox(
            Inches(l + indent), Inches(y), Inches(w - indent), Inches(spacing))
        txb.word_wrap = True
        tf = txb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = f"{m}  {text}"
        run.font.size = Pt(s)
        run.font.color.rgb = color
        y += spacing
    return y

def slide_header(slide, title, subtitle=None):
    """Темна смуга зверху з заголовком."""
    rect(slide, 0, 0, 13.33, 1.3, DARK_BLUE)
    txt(slide, title, 0.4, 0.1, 12.5, 0.75,
        size=28, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        txt(slide, subtitle, 0.4, 0.82, 12.5, 0.4,
            size=16, color=RGBColor(0xAD, 0xB5, 0xBD), align=PP_ALIGN.LEFT)
    # Акцентна лінія
    line = slide.shapes.add_shape(1,
        Inches(0), Inches(1.3), Inches(13.33), Inches(0.05))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT_BLUE
    line.line.fill.background()

def accent_card(slide, l, t, w, h, title, body_lines, icon="", title_size=17, body_size=14):
    """Картка з кольоровою лівою смугою."""
    rect(slide, l, t, 0.06, h, ACCENT_BLUE)
    card = rect(slide, l + 0.06, t, w - 0.06, h, RGBColor(0xF8, 0xF9, 0xFF))
    txt(slide, f"{icon}  {title}" if icon else title,
        l + 0.18, t + 0.08, w - 0.25, 0.35,
        size=title_size, bold=True, color=DARK_BLUE)
    y = t + 0.42
    for line in body_lines:
        txt(slide, line, l + 0.18, y, w - 0.28, 0.32,
            size=body_size, color=DARK_TEXT)
        y += 0.30

# ══════════════════════════════════════════════════════════════════════════════
# СЛАЙД 1 — ТИТУЛЬНИЙ
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, DARK_BLUE)

# Діагональний акцентний блок (декор)
dec = s.shapes.add_shape(1, Inches(8.5), Inches(0), Inches(5), Inches(7.5))
dec.fill.solid()
dec.fill.fore_color.rgb = RGBColor(0x0A, 0x1A, 0x3A)
dec.line.fill.background()

# Акцентна вертикальна лінія
rect(s, 8.48, 0, 0.07, 7.5, ACCENT_BLUE)

# Логотип/іконка (текстова)
txt(s, "🖥", 0.5, 0.5, 2, 1.5, size=60, color=ACCENT_BLUE)

txt(s, "ДЕРЖАВНИЙ УНІВЕРСИТЕТ", 0.5, 0.4, 8, 0.4,
    size=13, color=RGBColor(0xAD, 0xB5, 0xBD))
txt(s, "ІНФОРМАЦІЙНО-КОМУНІКАЦІЙНИХ ТЕХНОЛОГІЙ", 0.5, 0.72, 8, 0.4,
    size=13, color=RGBColor(0xAD, 0xB5, 0xBD))
txt(s, "Кафедра інженерії програмного забезпечення", 0.5, 1.0, 8, 0.35,
    size=12, italic=True, color=RGBColor(0x8A, 0x95, 0xAA))

# Розділова лінія
rect(s, 0.5, 1.45, 7.8, 0.04, ACCENT_BLUE)

txt(s, "КУРСОВИЙ ПРОЕКТ", 0.5, 1.6, 8, 0.5,
    size=16, color=RGBColor(0xAD, 0xB5, 0xBD))

txt(s, "Система управління інвентарем\nдля ІТ-департаменту", 0.5, 2.1, 8, 1.4,
    size=32, bold=True, color=WHITE)

txt(s, "Облік ліцензій, обладнання, серійних номерів\nз ролями користувачів та авторизацією на основі JWT",
    0.5, 3.5, 8, 0.9, size=17, italic=True,
    color=RGBColor(0xAD, 0xB5, 0xBD))

rect(s, 0.5, 4.55, 7.8, 0.04, RGBColor(0x2A, 0x3A, 0x5A))

txt(s, "Виконав:  студент гр. ТЦР-33  Скакун Михайло Миколайович",
    0.5, 4.7, 8, 0.4, size=15, color=WHITE)
txt(s, "Керівник:  Ніщеменко Дмитро Олександрович",
    0.5, 5.05, 8, 0.4, size=15, color=WHITE)
txt(s, "Київ — 2026", 0.5, 6.9, 8, 0.4,
    size=13, color=RGBColor(0x6C, 0x75, 0x7D))

# Права панель — технологічний стек
txt(s, "ТЕХНОЛОГІЧНИЙ СТЕК", 8.8, 0.5, 4.2, 0.4,
    size=13, bold=True, color=RGBColor(0xAD, 0xB5, 0xBD))
techs = [
    ("☕", "Java 17"),
    ("🍃", "Spring Boot 3.3.5"),
    ("🔐", "Spring Security + JWT"),
    ("🗄", "PostgreSQL 16"),
    ("📨", "RabbitMQ"),
    ("🐳", "Docker"),
    ("🧪", "Testcontainers"),
]
y = 1.1
for icon, name in techs:
    txt(s, f"{icon}  {name}", 8.8, y, 4.2, 0.45, size=16, color=WHITE)
    y += 0.72

# ══════════════════════════════════════════════════════════════════════════════
# СЛАЙД 2 — АКТУАЛЬНІСТЬ ТА МЕТА
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, WHITE)
slide_header(s, "Актуальність та мета роботи")

# Ліва колонка — проблема
rect(s, 0.3, 1.5, 6.0, 5.7, RGBColor(0xFF, 0xF3, 0xCD))
rect(s, 0.3, 1.5, 0.08, 5.7, RGBColor(0xFF, 0xC1, 0x07))
txt(s, "⚠  ПРОБЛЕМА", 0.55, 1.6, 5.5, 0.45, size=17, bold=True,
    color=RGBColor(0x85, 0x60, 0x04))
problems = [
    "Облік обладнання ведеться у Excel-таблицях",
    "Відсутній контроль терміну дії ліцензій",
    "Немає розмежування доступу за ролями",
    "Втрата інформації при звільненні співробітника",
    "Немає аудиту дій адміністраторів",
    "Ручний облік → помилки та дублювання даних",
]
add_bullet_box(s, problems, 0.55, 2.15, 5.6, 5.0,
               size=16, color=RGBColor(0x5A, 0x40, 0x00), spacing=0.57)

# Права колонка — рішення
rect(s, 6.8, 1.5, 6.2, 5.7, RGBColor(0xD1, 0xE7, 0xDD))
rect(s, 6.8, 1.5, 0.08, 5.7, GREEN)
txt(s, "✅  РІШЕННЯ", 7.05, 1.6, 5.8, 0.45, size=17, bold=True,
    color=RGBColor(0x0A, 0x53, 0x2E))
solutions = [
    "Централізована веб-система для ІТ-департаменту",
    "RBAC: чотири ролі (Admin/TeamLead/Worker/Auditor)",
    "JWT-авторизація для REST API",
    "Сповіщення про прострочені ліцензії (RabbitMQ)",
    "Повний аудит-лог усіх змін",
    "Excel-звіти по обладнанню та ліцензіях",
]
add_bullet_box(s, solutions, 7.05, 2.15, 5.8, 5.0,
               size=16, color=RGBColor(0x0A, 0x53, 0x2E), spacing=0.57)

txt(s, "Мета: розробити систему управління ІТ-інвентарем з повноцінною автентифікацією, розмежуванням доступу та автоматизованим обліком активів.",
    0.3, 7.0, 12.8, 0.45, size=14, italic=True, color=GRAY)

# ══════════════════════════════════════════════════════════════════════════════
# СЛАЙД 3 — АРХІТЕКТУРА СИСТЕМИ
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, WHITE)
slide_header(s, "Архітектура системи", "Тришарова монолітна архітектура + асинхронна черга")

layers = [
    ("🌐  Presentation Layer", ACCENT_BLUE,
     ["Thymeleaf (веб-інтерфейс)", "REST Controllers (/api/v1/...)", "Swagger UI / OpenAPI 3"]),
    ("⚙️  Business Logic Layer", RGBColor(0x19, 0x87, 0x54),
     ["Services (Equipment, License, User...)", "Spring Security + JWT Filter", "RabbitMQ Producer/Consumer"]),
    ("🗄  Data Access Layer", RGBColor(0xDC, 0x35, 0x45),
     ["JPA Repositories + Hibernate 6", "Flyway Migrations (V1–V4)", "PostgreSQL 16"]),
]

y = 1.5
for title, color, items in layers:
    rect(s, 0.5, y, 8.5, 1.55, RGBColor(0xF8, 0xF9, 0xFF))
    rect(s, 0.5, y, 0.12, 1.55, color)
    txt(s, title, 0.75, y + 0.08, 8, 0.45, size=18, bold=True, color=color)
    txt(s, "   ·  ".join(items), 0.75, y + 0.55, 8, 0.7, size=14, color=DARK_TEXT)
    if y < 4.5:
        txt(s, "↓", 4.5, y + 1.55, 0.5, 0.4, size=20, color=GRAY, align=PP_ALIGN.CENTER)
    y += 1.95

# Права панель — додаткові компоненти
txt(s, "Зовнішні компоненти", 9.4, 1.5, 3.7, 0.4,
    size=15, bold=True, color=DARK_BLUE)
ext = [
    ("🐳", "Docker Compose", "PostgreSQL + RabbitMQ"),
    ("🔑", "JJWT 0.12.6", "Access + Refresh tokens"),
    ("📊", "Apache POI 5.3", "Excel-звіти (.xlsx)"),
    ("🧪", "Testcontainers", "Інтеграційні тести"),
]
y = 2.0
for icon, name, desc in ext:
    rect(s, 9.4, y, 3.7, 1.1, RGBColor(0xF0, 0xF4, 0xFF))
    txt(s, f"{icon} {name}", 9.6, y + 0.08, 3.4, 0.38, size=15, bold=True, color=DARK_BLUE)
    txt(s, desc, 9.6, y + 0.45, 3.4, 0.38, size=13, color=GRAY)
    y += 1.2

# ══════════════════════════════════════════════════════════════════════════════
# СЛАЙД 4 — РОЛІ ТА ПРАВА ДОСТУПУ (RBAC)
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, WHITE)
slide_header(s, "Модель розмежування доступу (RBAC)", "Role-Based Access Control — 4 ролі, матриця прав")

roles = [
    ("👑  ADMIN",    RGBColor(0xDC,0x35,0x45), [
        "Повний доступ до всіх ресурсів",
        "Управління користувачами та об'єктами",
        "Перегляд журналу аудиту",
        "Генерація Excel-звітів",
    ]),
    ("🏆  TEAM LEAD", RGBColor(0x0D,0x6E,0xFD), [
        "Управління командою свого об'єкта",
        "Додавання/редагування обладнання та ліцензій",
        "Перегляд свого об'єкта",
        "Генерація звітів по своєму об'єкту",
    ]),
    ("👷  WORKER",   RGBColor(0x19,0x87,0x54), [
        "Перегляд свого інвентарю",
        "Перегляд призначених ліцензій",
        "Особистий дашборд",
        "Без доступу до адмін-панелі",
    ]),
    ("🔍  AUDITOR",  RGBColor(0x6C,0x75,0x7D), [
        "Read-only доступ до всіх даних",
        "Перегляд журналу аудиту",
        "Без права на зміни",
        "Звіти та статистика",
    ]),
]

x_positions = [0.3, 3.45, 6.6, 9.75]
for i, (role, color, perms) in enumerate(roles):
    x = x_positions[i]
    rect(s, x, 1.45, 3.0, 0.55, color)
    txt(s, role, x + 0.1, 1.5, 2.8, 0.45, size=16, bold=True, color=WHITE)
    rect(s, x, 2.0, 3.0, 4.8, RGBColor(0xF8,0xF9,0xFF))
    y = 2.1
    for perm in perms:
        txt(s, f"✓  {perm}", x + 0.1, y, 2.8, 0.55, size=13, color=DARK_TEXT)
        y += 0.58
    rect(s, x, 2.0, 0.06, 4.8, color)

txt(s, "Реалізація: @PreAuthorize(\"hasRole('ADMIN')\") + InventoryPermissionEvaluator (hasPermission) на рівні сервісів",
    0.3, 7.05, 12.8, 0.4, size=13, italic=True, color=GRAY)

# ══════════════════════════════════════════════════════════════════════════════
# СЛАЙД 5 — JWT-АВТОРИЗАЦІЯ
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, WHITE)
slide_header(s, "Система автентифікації — JWT", "Два SecurityFilterChain: stateless API + session Web")

# Схема потоку
steps = [
    ("1", "POST /api/v1/auth/login\n{username, password}", ACCENT_BLUE),
    ("2", "BCrypt verify\npassword hash", RGBColor(0x6F,0x42,0xC1)),
    ("3", "generateAccessToken\n(15 хв, HS256)", GREEN),
    ("4", "Refresh Token\n(30 днів, DB)", RGBColor(0xFF,0x80,0x00)),
    ("5", "Authorization: Bearer\n<token>", DARK_BLUE),
]

x = 0.25
for i, (num, label, color) in enumerate(steps):
    rect(s, x, 1.6, 2.3, 1.4, color)
    txt(s, num, x + 0.08, 1.65, 0.4, 0.5, size=22, bold=True, color=WHITE)
    txt(s, label, x + 0.08, 2.0, 2.1, 0.9, size=13, color=WHITE)
    if i < 4:
        txt(s, "→", x + 2.3, 2.1, 0.3, 0.5, size=22, bold=True, color=GRAY, align=PP_ALIGN.CENTER)
    x += 2.62

# Дві колонки — FilterChain
txt(s, "FilterChain 1 — REST API  (/api/**)", 0.5, 3.3, 6.0, 0.4,
    size=16, bold=True, color=ACCENT_BLUE)
fc1 = [
    "@Order(1)  — обробляється першим",
    "CSRF disabled, SessionPolicy = STATELESS",
    "JwtAuthenticationFilter перехоплює кожен запит",
    "Витягує userId з токена → SecurityContext",
    "Ендпоінти: /api/v1/equipment, /api/v1/licenses ...",
]
add_bullet_box(s, fc1, 0.5, 3.75, 6.0, 3.5, size=14, color=DARK_TEXT, spacing=0.52)

txt(s, "FilterChain 2 — Web / Thymeleaf  (/**)", 6.9, 3.3, 6.1, 0.4,
    size=16, bold=True, color=DARK_BLUE)
fc2 = [
    "@Order(2) — обробляє сесійний веб-трафік",
    "CSRF enabled (th:action додає токен)",
    "Form-login: /login → defaultSuccessUrl=/dashboard",
    "Session-based auth через JSESSIONID cookie",
    "Logout: POST /web/logout → /login?logout=true",
]
add_bullet_box(s, fc2, 6.9, 3.75, 6.1, 3.5, size=14, color=DARK_TEXT, spacing=0.52)

rect(s, 6.7, 3.25, 0.05, 4.0, RGBColor(0xDE, 0xE2, 0xE6))

# ══════════════════════════════════════════════════════════════════════════════
# СЛАЙД 6 — ШАБЛОНИ ПРОЕКТУВАННЯ
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, WHITE)
slide_header(s, "Шаблони проектування (GoF)", "7 патернів, реалізованих у проекті")

patterns = [
    ("🏭", "Factory Method", "EquipmentFactory",
     "forType(type) повертає ConcreteFactory залежно від типу обладнання"),
    ("📋", "Specification", "EquipmentSpecification",
     "Комбінує фільтри (siteId, status, type, q) через Specification<T>.and()"),
    ("👁", "Observer", "AuditLogListener",
     "@EventListener реагує на EntityChangedEvent і зберігає аудит (REQUIRES_NEW)"),
    ("🔀", "Strategy", "NotificationStrategy",
     "Email/InApp/Telegram — єдиний інтерфейс, динамічний вибір стратегії"),
    ("🔨", "Builder", "EquipmentReportRequest",
     "Lombok @Builder для конструювання складних DTO запитів на звіти"),
    ("🎨", "Decorator", "LicenseKeyAttributeConverter",
     "@Convert AES-256-GCM шифрує ліцензійні ключі прозоро для бізнес-логіки"),
    ("☝", "Singleton", "JwtService / ReportStore",
     "Spring singleton beans; ReportStore — ConcurrentHashMap для in-memory зберігання"),
]

cols = 4
x_pos = [0.3, 3.42, 6.54, 9.66]
y_pos = [1.5, 3.45]
for i, (icon, name, cls, desc) in enumerate(patterns):
    col = i % cols
    row = i // cols
    if col == 3 and row == 1:
        # Центруємо останній патерн (7-й) у 4-й колонці 2-го рядку
        pass
    x = x_pos[col]
    y = y_pos[row]
    rect(s, x, y, 3.0, 1.75, RGBColor(0xF0, 0xF4, 0xFF))
    rect(s, x, y, 0.07, 1.75, ACCENT_BLUE)
    txt(s, f"{icon}  {name}", x + 0.18, y + 0.08, 2.75, 0.38,
        size=15, bold=True, color=DARK_BLUE)
    txt(s, cls, x + 0.18, y + 0.45, 2.75, 0.3,
        size=12, italic=True, color=ACCENT_BLUE)
    txt(s, desc, x + 0.18, y + 0.78, 2.75, 0.82,
        size=11.5, color=DARK_TEXT)

# ══════════════════════════════════════════════════════════════════════════════
# СЛАЙД 7 — БАЗА ДАНИХ
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, WHITE)
slide_header(s, "Модель бази даних", "PostgreSQL 16 · 7 таблиць · Flyway міграції V1–V4")

tables_info = [
    ("sites",          "Об'єкти (відділення)",         "id, name, address, city, team_lead_id"),
    ("users",          "Користувачі системи",           "id, username, email, role, site_id, active"),
    ("equipment",      "Обладнання (активи)",           "id, inventory_number, type, status, assigned_user_id"),
    ("licenses",       "Програмні ліцензії",            "id, software_name, license_key (AES), seats_total/used"),
    ("audit_logs",     "Журнал аудиту (JSONB payload)", "id, actor_user_id, action, entity_type, entity_id"),
    ("refresh_tokens", "Refresh-токени (ротація)",      "id, user_id, token, expires_at, revoked"),
    ("notifications",  "Сповіщення (RabbitMQ)",         "id, recipient_user_id, type, status, payload"),
]

y = 1.55
colors = [ACCENT_BLUE, GREEN, RGBColor(0xDC,0x35,0x45),
          RGBColor(0xFF,0x80,0x00), RGBColor(0x6F,0x42,0xC1),
          GRAY, RGBColor(0x0D,0xCA,0xF0)]
for i, (table, desc, cols_info) in enumerate(tables_info):
    c = colors[i]
    rect(s, 0.3, y, 2.5, 0.72, c)
    txt(s, table, 0.4, y + 0.12, 2.3, 0.45, size=15, bold=True, color=WHITE)
    rect(s, 2.8, y, 10.2, 0.72, RGBColor(0xF8, 0xF9, 0xFF))
    txt(s, desc, 2.95, y + 0.04, 4.5, 0.3, size=14, bold=True, color=DARK_BLUE)
    txt(s, cols_info, 2.95, y + 0.38, 9.8, 0.28, size=12, italic=True, color=GRAY)
    y += 0.8

txt(s, "Особливості: UUID первинні ключі · CHECK constraints замість PostgreSQL ENUM (Hibernate 6 сумісність) · JSONB для payload · AES-256-GCM для license_key",
    0.3, 7.1, 12.8, 0.35, size=12, italic=True, color=GRAY)

# ══════════════════════════════════════════════════════════════════════════════
# СЛАЙД 8 — ДІАГРАМА КЛАСІВ
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, WHITE)
slide_header(s, "Діаграма класів", "UML — ключові пакети та зв'язки між ними")

if os.path.exists(UML_CLASS):
    pic = s.shapes.add_picture(UML_CLASS, Inches(0.2), Inches(1.4), Inches(12.9), Inches(5.9))
else:
    rect(s, 0.2, 1.4, 12.9, 5.9, RGBColor(0xEE, 0xF0, 0xF2))
    txt(s, "[Діаграма класів — IT-Inventory Class Diagram.png]",
        0.5, 4.0, 12.3, 0.6, size=18, color=GRAY, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# СЛАЙД 9 — SEQUENCE ДІАГРАМА JWT
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, WHITE)
slide_header(s, "Sequence-діаграма: JWT авторизація", "POST /api/v1/auth/login → Access Token + Refresh Token")

if os.path.exists(UML_SEQ):
    pic = s.shapes.add_picture(UML_SEQ, Inches(0.3), Inches(1.4), Inches(12.7), Inches(5.9))
else:
    rect(s, 0.3, 1.4, 12.7, 5.9, RGBColor(0xEE, 0xF0, 0xF2))
    txt(s, "[Sequence Login Diagram — Login Sequence.png]",
        0.5, 4.0, 12.3, 0.6, size=18, color=GRAY, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# СЛАЙД 10 — REST API
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, WHITE)
slide_header(s, "REST API — ключові ендпоінти", "Prefix: /api/v1/ · Формат помилок: RFC 7807 ProblemDetail")

endpoints = [
    ("POST",   "/auth/login",              "—",                   "Отримати access + refresh токени"),
    ("POST",   "/auth/refresh",            "—",                   "Оновити access токен (ротація refresh)"),
    ("GET",    "/equipment",               "ADMIN/LEAD/WORKER",   "Список обладнання (фільтр, пагінація)"),
    ("POST",   "/equipment",               "ADMIN/LEAD",          "Додати нове обладнання"),
    ("POST",   "/equipment/{id}/assign",   "ADMIN/LEAD",          "Призначити обладнання користувачу"),
    ("GET",    "/licenses",                "ADMIN/LEAD/WORKER",   "Список ліцензій"),
    ("POST",   "/licenses",                "ADMIN/LEAD",          "Додати нову ліцензію"),
    ("GET",    "/users",                   "ADMIN/AUDITOR",       "Список користувачів"),
    ("POST",   "/users",                   "ADMIN",               "Створити користувача"),
    ("GET",    "/audit-log",               "ADMIN/AUDITOR/LEAD",  "Журнал аудиту (фільтр)"),
    ("GET",    "/reports/inventory.xlsx",  "ADMIN/AUDITOR/LEAD",  "Запит на генерацію Excel-звіту"),
    ("GET",    "/reports/{id}/download",   "Authenticated",       "Завантажити готовий звіт"),
]

method_colors = {
    "GET":  RGBColor(0x19, 0x87, 0x54),
    "POST": ACCENT_BLUE,
    "PUT":  RGBColor(0xFF, 0x80, 0x00),
    "DELETE": RGBColor(0xDC, 0x35, 0x45),
}

# Заголовок таблиці
rect(s, 0.25, 1.45, 12.85, 0.45, DARK_BLUE)
for x, w, t in [(0.3,1.0,"Метод"),(1.35,3.5,"Шлях"),(4.9,2.8,"Роль"),(7.75,5.3,"Опис")]:
    txt(s, t, x, 1.5, w, 0.35, size=13, bold=True, color=WHITE)

y = 1.9
for i, (method, path, role, desc) in enumerate(endpoints):
    bg_c = RGBColor(0xF8,0xF9,0xFF) if i % 2 == 0 else WHITE
    rect(s, 0.25, y, 12.85, 0.44, bg_c)
    mc = method_colors.get(method, GRAY)
    rect(s, 0.25, y+0.05, 1.0, 0.34, mc)
    txt(s, method, 0.28, y+0.07, 0.94, 0.3, size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(s, path,   1.35, y+0.06, 3.5, 0.32, size=11.5, color=DARK_BLUE)
    txt(s, role,   4.9,  y+0.06, 2.8, 0.32, size=11, color=RGBColor(0x6F,0x42,0xC1))
    txt(s, desc,   7.75, y+0.06, 5.3, 0.32, size=11.5, color=DARK_TEXT)
    y += 0.44

# ══════════════════════════════════════════════════════════════════════════════
# СЛАЙД 11 — ВЕББІНТЕРФЕЙС (скриншоти — плейсхолдери)
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, WHITE)
slide_header(s, "Веб-інтерфейс — Thymeleaf + Bootstrap 5", "Адаптивний інтерфейс з рольовою видимістю елементів")

screens = [
    ("Сторінка входу", "login.png"),
    ("Дашборд адміна", "dashboard_admin.png"),
    ("Список обладнання", "equipment_list.png"),
    ("Журнал аудиту", "audit_log.png"),
]

x_pos = [0.25, 6.75]
y_pos = [1.45, 4.55]
for i, (title, fname) in enumerate(screens):
    x = x_pos[i % 2]
    y = y_pos[i // 2]
    path = f"/Users/mihailskakun/IdeaProjects/kursova_java/docs/screenshots/{fname}"
    if os.path.exists(path):
        s.shapes.add_picture(path, Inches(x), Inches(y), Inches(6.3), Inches(2.85))
    else:
        rect(s, x, y, 6.3, 2.85, RGBColor(0xE9, 0xEC, 0xEF))
        txt(s, f"📱 {title}", x+0.15, y+1.1, 6.0, 0.55,
            size=16, color=GRAY, align=PP_ALIGN.CENTER)
        txt(s, f"[screenshots/{fname}]", x+0.15, y+1.7, 6.0, 0.4,
            size=12, italic=True, color=GRAY, align=PP_ALIGN.CENTER)

txt(s, "Особливості UI: sec:authorize на <li> елементах · th:action для CSRF · Пагінація · Фільтри · Bootstrap 5 модальні вікна",
    0.25, 7.1, 12.85, 0.35, size=12, italic=True, color=GRAY)

# ══════════════════════════════════════════════════════════════════════════════
# СЛАЙД 12 — ТЕСТУВАННЯ
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, WHITE)
slide_header(s, "Тестування", "Unit tests (Mockito) + Integration tests (Testcontainers) + JaCoCo")

# Ліва колонка — unit тести
txt(s, "🧪  Юніт-тести (Mockito)", 0.3, 1.55, 6.0, 0.45,
    size=17, bold=True, color=DARK_BLUE)
units = [
    "JwtServiceTest — генерація/валідація токена",
    "EquipmentServiceTest — create/assign бізнес-правила",
    "  conflictException при дублікаті serialNumber",
    "  businessRuleException при різних siteId",
    "LicenseServiceTest — перевірка seatsUsed/Total",
    "AuditLogListenerTest — Observer, REQUIRES_NEW TX",
]
add_bullet_box(s, units, 0.3, 2.05, 6.2, 5.0, size=14, spacing=0.55)

# Права колонка — integration тести
txt(s, "🔗  Інтеграційні тести (Testcontainers)", 6.9, 1.55, 6.1, 0.45,
    size=17, bold=True, color=DARK_BLUE)
integr = [
    "AbstractIntegrationTest — PostgreSQL + RabbitMQ у Docker",
    "AuthIntegrationTest (6 тестів):",
    "  login valid/invalid credentials → 200/401",
    "  token refresh + ротація refresh token",
    "  protected endpoint без/з токеном → 401/200",
    "EquipmentIntegrationTest — CRUD + аудит + 403",
]
add_bullet_box(s, integr, 6.9, 2.05, 6.1, 5.0, size=14, spacing=0.55)

rect(s, 6.7, 1.5, 0.05, 5.7, RGBColor(0xDE, 0xE2, 0xE6))

# JaCoCo підсумок
rect(s, 0.3, 6.65, 12.73, 0.75, RGBColor(0xD1, 0xE7, 0xDD))
txt(s, "✅  JaCoCo: мінімум покриття 60% (налаштовано у pom.xml). "
       "Звіт: target/site/jacoco/index.html. "
       "Всього тестів: 31 (20 unit + 11 integration).",
    0.5, 6.75, 12.4, 0.5, size=14, color=RGBColor(0x0A, 0x53, 0x2E))

# ══════════════════════════════════════════════════════════════════════════════
# СЛАЙД 13 — КОНТЕЙНЕРИЗАЦІЯ
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, WHITE)
slide_header(s, "Контейнеризація — Docker", "Multi-stage Dockerfile + docker-compose.yml")

# Multi-stage build
txt(s, "Multi-stage Dockerfile", 0.3, 1.55, 5.8, 0.4, size=17, bold=True, color=DARK_BLUE)
dockerfile_lines = [
    "# Stage 1 — збірка",
    "FROM eclipse-temurin:17-jdk-alpine AS build",
    "COPY . .",
    "RUN mvn package -DskipTests",
    "",
    "# Stage 2 — runtime (мінімальний образ)",
    "FROM eclipse-temurin:17-jre-alpine",
    "COPY --from=build target/*.jar app.jar",
    'ENTRYPOINT ["java","-jar","/app.jar"]',
]
rect(s, 0.3, 2.0, 5.9, 3.5, RGBColor(0x1E, 0x1E, 0x2E))
y = 2.1
for line in dockerfile_lines:
    color = RGBColor(0x89, 0xDD, 0xFF) if line.startswith("#") else \
            RGBColor(0xC3, 0xE8, 0x8D) if line.startswith("FROM") else \
            RGBColor(0xFF, 0xFF, 0xFF)
    txt(s, line, 0.45, y, 5.6, 0.28, size=11.5, color=color)
    y += 0.29

# Docker Compose
txt(s, "docker-compose.yml — сервіси", 6.6, 1.55, 6.5, 0.4, size=17, bold=True, color=DARK_BLUE)
services = [
    ("🗄  inventory_postgres", "postgres:16-alpine",
     "Port 5432 · DB: inventory · Volume: pgdata"),
    ("📨  inventory_rabbitmq", "rabbitmq:3-management-alpine",
     "AMQP 5672 · UI 15672 · healthcheck"),
    ("☕  inventory_app", "inventory-system:latest",
     "Port 8080 · залежить від postgres + rabbit · profiles: prod"),
]
y = 2.0
for name, image, desc in services:
    rect(s, 6.6, y, 6.5, 1.3, RGBColor(0xF0, 0xF4, 0xFF))
    rect(s, 6.6, y, 0.08, 1.3, ACCENT_BLUE)
    txt(s, name,  6.75, y+0.08, 6.2, 0.38, size=15, bold=True, color=DARK_BLUE)
    txt(s, image, 6.75, y+0.46, 6.2, 0.28, size=12, italic=True, color=ACCENT_BLUE)
    txt(s, desc,  6.75, y+0.78, 6.2, 0.38, size=12, color=DARK_TEXT)
    y += 1.42

txt(s, "Запуск: docker compose up -d     |     Зупинка: docker compose down",
    0.3, 5.65, 12.8, 0.4, size=14, color=DARK_BLUE,
    align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# СЛАЙД 14 — ВИСНОВКИ
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, DARK_BLUE)

rect(s, 0, 0, 13.33, 0.08, ACCENT_BLUE)

txt(s, "ВИСНОВКИ", 0.5, 0.3, 12.3, 0.65,
    size=32, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

results = [
    ("✅", "Розроблено повнофункціональну систему управління ІТ-інвентарем"),
    ("✅", "Реалізовано RBAC з 4 ролями та JWT-авторизацію (access + refresh token)"),
    ("✅", "Застосовано 7 шаблонів проектування GoF у реальному коді"),
    ("✅", "Реалізовано асинхронну генерацію Excel-звітів через RabbitMQ"),
    ("✅", "Написано 31 тест (unit + integration) з покриттям ≥ 60% (JaCoCo)"),
    ("✅", "Проект контейнеризовано через Docker + docker-compose"),
    ("✅", "Документацію API згенеровано через SpringDoc / Swagger UI"),
]

y = 1.2
for icon, text in results:
    rect(s, 0.4, y, 12.5, 0.68, RGBColor(0x0A, 0x1A, 0x3A))
    txt(s, icon, 0.55, y+0.1, 0.5, 0.45, size=20, color=GREEN)
    txt(s, text, 1.15, y+0.13, 11.7, 0.42, size=16, color=WHITE)
    y += 0.76

txt(s, "Перспективи розвитку: мобільний застосунок · QR-коди на обладнання · Telegram-бот для сповіщень · SSO (OAuth2)",
    0.4, 6.75, 12.5, 0.45, size=13, italic=True,
    color=RGBColor(0xAD, 0xB5, 0xBD), align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# СЛАЙД 15 — ДЯКУЮ ЗА УВАГУ
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, DARK_BLUE)

rect(s, 0, 0, 13.33, 0.08, ACCENT_BLUE)
rect(s, 0, 7.42, 13.33, 0.08, ACCENT_BLUE)

txt(s, "🖥", 5.5, 0.8, 2.3, 1.8, size=80, color=ACCENT_BLUE, align=PP_ALIGN.CENTER)

txt(s, "Дякую за увагу!", 0.5, 2.6, 12.3, 1.0,
    size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

txt(s, "Система управління інвентарем для ІТ-департаменту", 0.5, 3.65, 12.3, 0.6,
    size=20, italic=True, color=RGBColor(0xAD, 0xB5, 0xBD), align=PP_ALIGN.CENTER)

rect(s, 3.5, 4.4, 6.33, 0.05, RGBColor(0x2A, 0x3A, 0x5A))

info = [
    ("👤", "Скакун Михайло Миколайович"),
    ("🎓", "Група ТЦР-33  ·  ДУІКТ"),
    ("👨‍🏫", "Керівник: Ніщеменко Д. О."),
    ("💻", "github.com/TimeDustik/kursova_java"),
]
y = 4.65
for icon, text in info:
    txt(s, f"{icon}  {text}", 0.5, y, 12.3, 0.45,
        size=17, color=WHITE, align=PP_ALIGN.CENTER)
    y += 0.5

txt(s, "Київ — 2026", 0.5, 7.0, 12.3, 0.38,
    size=14, color=RGBColor(0x6C, 0x75, 0x7D), align=PP_ALIGN.CENTER)

# ─── Зберегти ─────────────────────────────────────────────────────────────────
prs.save(OUTPUT)

size_kb = os.path.getsize(OUTPUT) // 1024
print("=" * 60)
print(f"Презентацію збережено: {OUTPUT}")
print(f"Розмір файлу:  {size_kb} КБ")
print(f"Кількість слайдів: {len(prs.slides)}")
print("=" * 60)
print("Слайди:")
titles = [
    "1. Титульний",
    "2. Актуальність та мета",
    "3. Архітектура системи",
    "4. Ролі та права доступу (RBAC)",
    "5. JWT-авторизація",
    "6. Шаблони проектування (GoF)",
    "7. Модель бази даних",
    "8. Діаграма класів",
    "9. Sequence-діаграма JWT",
    "10. REST API ендпоінти",
    "11. Веб-інтерфейс (Thymeleaf)",
    "12. Тестування (JaCoCo)",
    "13. Контейнеризація (Docker)",
    "14. Висновки",
    "15. Дякую за увагу",
]
for t in titles:
    print(f"  {t}")
